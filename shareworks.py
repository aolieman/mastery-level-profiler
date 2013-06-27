#!/usr/bin/env python
from __future__ import division
import os, itertools
import layout_scanner
from bs4 import BeautifulSoup
from docx import opendocx, getdocumenttext
from pptx import Presentation

def readPortfolios(portfolios_dir, verbose=0):
    # this function assumes the indir contains files only
    filelist = os.listdir(portfolios_dir)
    dictlist = []
    for f in filelist:
        fextstr = f.split(".")[-1]
        print "\n\n", fextstr, f
        ## PDF
        if fextstr == "pdf":
            pdf_dict = convPdf(portfolios_dir, f)
            dictlist.append(pdf_dict)
            if verbose > 0: print "\n", pdf_dict
        ## PPTX
        elif fextstr == "pptx":
            pptx_dict = convPptx(portfolios_dir, f)
            dictlist.append(pptx_dict)
            if verbose > 0: print "\n", pptx_dict
        ## DOCX
        elif fextstr == "docx":
            docx_dict = convDocx(portfolios_dir, f)
            dictlist.append(docx_dict)
            if verbose > 0: print "\n", docx_dict
        ## HTML
        elif fextstr in ("html", "htm"):
            html_dict = convHtml(portfolios_dir, f)
            dictlist.append(html_dict)
            if verbose > 0: print "\n", html_dict
        else: raise Exception("Check if this filetype can be read!")

    return dictlist

def convPdf(portfolios_dir, f):
    pdf_filepath = os.path.join(portfolios_dir, f)
    toc=layout_scanner.get_toc(pdf_filepath)
    print "TOC -- not currently used for dict"
    for e in toc: print e
    pages=layout_scanner.get_pages(pdf_filepath)
    doc_dict = {'doctype': "report", 'origin': "shareworks"}
    doc_dict['_id'] = f
    doc_dict['student_email'] = email_from_fname(f)
    doc_dict['content'] = []
    for (count, p) in enumerate(pages):
        if len(p) > 5: # when is a page non-empty?
            doc_dict['content'].append({'header': "Page %s" % count, 'text': p.decode('utf-8')})
    if len(doc_dict['content']) == 0:
        doc_dict['content'] = [{'text': 'FrozenCutlery', 'txt_ann': None}]
    
    return doc_dict

def convPptx(portfolios_dir, f):
    pptx_filepath = os.path.join(portfolios_dir, f)
    prs = Presentation(pptx_filepath)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_textframe:
                continue
            for paragraph in shape.textframe.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    all_slides = "\n".join(text_runs)
    
    doc_dict = {'doctype': "slides", 'origin': "shareworks"}
    doc_dict['_id'] = f
    doc_dict['student_email'] = email_from_fname(f)
    doc_dict['content'] = [{'header': "All slides", 'text': all_slides}]
    if len(doc_dict['content'][0]['text']) == 0:
        doc_dict['content'] = [{'text': 'FrozenCutlery', 'txt_ann': None}]
    
    return doc_dict

def convDocx(portfolios_dir, f):
    docx_filepath = os.path.join(portfolios_dir, f)
    docu = opendocx(docx_filepath)
    paratextlist = getdocumenttext(docu)
    
    doc_dict = {'doctype': "report", 'origin': "shareworks"}
    doc_dict['_id'] = f
    doc_dict['student_email'] = email_from_fname(f)
    doc_dict['content'] = []
    for pair in pairwise(paratextlist):
        if len(pair[0]) / len(pair[1]) < 1.5: #assume header is not much longer than text
            doc_dict['content'].append({'header': pair[0], 'text': "\n\n".join(pair)})
    if len(doc_dict['content']) == 0:
        doc_dict['content'] = [{'text': 'FrozenCutlery', 'txt_ann': None}]
    
    return doc_dict

def convHtml(portfolios_dir, f):
    html_filepath = os.path.join(portfolios_dir, f)
    html_file = open(html_filepath, 'r')
    soup = BeautifulSoup(html_file, "lxml")
    html_file.close()
    
    doc_dict = {'doctype': "posts", 'origin': "shareworks"}
    doc_dict['_id'] = f
    doc_dict['student_email'] = soup.head['id'].split(" ")
    doc_dict['title'] = soup.find("h1").text
    doc_dict['content'] = []
    posts = soup.find_all('div', 'sitepagenote')
    for post in posts:
        doc_dict['content'].append({'header': post.h3.text, 'text': post.text})
    if len(doc_dict['content']) == 0:
        doc_dict['content'] = [{'text': 'FrozenCutlery', 'txt_ann': None}]
    
    return doc_dict

def pairwise(iterable):
    "s -> (s0,s1), (s1,s2), (s2, s3), ..."
    a, b = itertools.tee(iterable)
    next(b, None)
    return itertools.izip(a, b)

def email_from_fname(f):
    if len(f.split("--")) == 2:
        return [f.split("--")[0]]
    else: return f.split("--")[:-1]

if __name__ == '__main__' :
    portfolios_path = "../Phase B/sw_portfolios"
    websites_path = "../Phase B/websites"

    readPortfolios(portfolios_path)
